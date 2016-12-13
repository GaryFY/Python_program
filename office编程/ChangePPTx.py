from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]#选取一个位置
slide = prs.slides.add_slide(title_slide_layout)#在位置上增加一个幻灯片
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello,World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')